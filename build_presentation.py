from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = "apresentacao_rov_control.pptx"
W, H = Inches(13.333), Inches(7.5)
BLACK = RGBColor(17, 17, 17)
GRAY = RGBColor(95, 99, 104)
LIGHT = RGBColor(242, 244, 247)
MID = RGBColor(211, 215, 220)
BLUE = RGBColor(22, 89, 160)
CYAN = RGBColor(0, 134, 179)
GREEN = RGBColor(38, 126, 79)
ORANGE = RGBColor(196, 101, 24)
RED = RGBColor(177, 45, 45)
WHITE = RGBColor(255, 255, 255)
FONT = "Open Sans"


def rgb(hexstr):
    return RGBColor.from_string(hexstr.replace("#", ""))


def set_run(run, size=18, bold=False, color=BLACK, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def textbox(slide, x, y, w, h, text="", size=18, bold=False, color=BLACK,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run(r, size, bold, color)
    return shape


def richtext(slide, x, y, w, h, paragraphs, margin=0.08):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    for i, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = spec["text"]
        p.level = spec.get("level", 0)
        p.space_after = Pt(spec.get("after", 7))
        p.line_spacing = spec.get("line", 1.05)
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        p.font.name = FONT
        p.font.size = Pt(spec.get("size", 17))
        p.font.bold = spec.get("bold", False)
        p.font.color.rgb = spec.get("color", BLACK)
        if spec.get("bullet", False):
            p.text = "• " + p.text
        for run in p.runs:
            set_run(run, spec.get("size", 17), spec.get("bold", False),
                    spec.get("color", BLACK))
    return shape


def box(slide, x, y, w, h, title, body="", fill=LIGHT, line=MID,
        title_color=BLACK, body_color=GRAY):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(1.2)
    textbox(slide, x + .16, y + .13, w - .32, .35, title, 16, True, title_color)
    if body:
        textbox(slide, x + .16, y + .56, w - .32, h - .66, body, 13.5, False, body_color)
    return sh


def line(slide, x1, y1, x2, y2, color=BLACK, width=1.5, arrow=False, dash=False):
    sh = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
    sh.line.color.rgb = color
    sh.line.width = Pt(width)
    if arrow:
        sh.line.end_arrowhead = True
    if dash:
        sh.line.dash_style = 2
    return sh


def pill(slide, x, y, w, text, fill, color=WHITE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(.34))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    textbox(slide, x, y + .01, w, .28, text, 10.5, True, color,
            PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 0)


def title(slide, n, heading, kicker=None):
    textbox(slide, .55, .26, .55, .32, f"{n:02d}", 12, True, BLUE)
    textbox(slide, 1.05, .2, 11.7, .56, heading, 26, True)
    if kicker:
        textbox(slide, 1.07, .72, 11.4, .34, kicker, 12.5, False, GRAY)
    line(slide, .55, 1.08, 12.78, 1.08, MID, .8)


def footer(slide, source):
    textbox(slide, .58, 7.15, 12.2, .18, source, 7.5, False, GRAY)


def add_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid(); bg.fore_color.rgb = WHITE
    return s


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    prs.core_properties.title = "Controle distribuído de ROVs"
    prs.core_properties.subject = "Arquitetura, protocolo QUIC-lite, segurança e tolerância a falhas"

    # 1 — problema
    s = add_slide(prs)
    textbox(s, .65, .55, 7.3, .38, "SISTEMAS DISTRIBUÍDOS · ROV CONTROL", 12, True, BLUE)
    textbox(s, .65, 1.08, 7.5, 1.35, "Controle remoto,\nsem perder o controle", 34, True)
    textbox(s, .7, 2.75, 6.5, 1.25,
            "Simular os componentes de um sistema completo de controle de veículos operados remotamente — com baixa latência, autenticação e tolerância a falhas.",
            18, False, GRAY)
    box(s, 8.15, .75, 4.25, 1.32, "ORIGEM DO PROBLEMA",
        "A ideia nasce da experiência de Fernando na OceanPact, desenvolvendo funcionalidades ligadas ao controle desses dispositivos.",
        LIGHT, MID, BLUE, BLACK)
    box(s, 8.15, 2.3, 4.25, 1.55, "O QUE É UM ROV?",
        "Um robô móvel operado à distância. Neste projeto: um veículo submarino para inspeção, pesquisa e apoio à exploração offshore de petróleo.",
        LIGHT, MID, CYAN, BLACK)
    box(s, 8.15, 4.1, 4.25, 1.78, "A ANALOGIA COM MARTE",
        "Curiosity e Perseverance são rovers robóticos, mas não ROVs em teleoperação contínua: a distância impede controle em tempo real. A Terra envia planos; a autonomia local executa e evita obstáculos.",
        LIGHT, MID, ORANGE, BLACK)
    pill(s, .7, 5.25, 2.0, "PILOTO", BLUE)
    line(s, 2.7, 5.42, 4.18, 5.42, BLUE, 2.3, True)
    pill(s, 4.2, 5.25, 2.0, "RELAY", CYAN)
    line(s, 6.2, 5.42, 7.68, 5.42, CYAN, 2.3, True)
    pill(s, 7.7, 5.25, 2.0, "ROV", GREEN)
    textbox(s, .7, 6.18, 9.05, .54,
            "DESAFIO  →  reproduzir comunicação, concorrência, segurança e recuperação de falhas em uma única simulação.",
            15, True)
    footer(s, "Fontes: README.md; NASA/JPL — How Perseverance Drives on Mars; NASA — AI-Planned Drive (teleoperação em tempo real é inviável).")

    # 2 — componentes e escopo
    s = add_slide(prs); title(s, 2, "O sistema: três papéis, cinco nós lógicos",
                              "Todos iniciam tráfego em direção ao relay conhecido; isso contorna NAT/firewall.")
    box(s, .65, 1.4, 3.65, 1.55, "PILOTO",
        "Autentica-se, solicita posse e envia comandos. Um piloto pode escolher entre vários ROVs — nesta demo, cada instância aponta para um alvo.",
        LIGHT, MID, BLUE, BLACK)
    box(s, 4.82, 1.4, 3.65, 1.55, "RELAY",
        "Ponto de encontro e roteador: mantém sessões, autentica, arbitra exclusão mútua e encaminha comandos/telemetria.",
        LIGHT, MID, CYAN, BLACK)
    box(s, 8.98, 1.4, 3.65, 1.55, "ROV",
        "Registra-se, aplica comandos e publica bateria, profundidade, temperatura e potência do thruster.",
        LIGHT, MID, GREEN, BLACK)
    line(s, 4.3, 2.17, 4.8, 2.17, BLUE, 2, True)
    line(s, 8.47, 2.17, 8.97, 2.17, CYAN, 2, True)
    textbox(s, .7, 3.3, 5.7, .3, "REQUISITOS COBERTOS", 14, True, GREEN)
    richtext(s, .68, 3.7, 5.85, 2.72, [
        {"text": "Vários pilotos e ROVs por relay; 1 piloto por ROV", "bullet": True, "size": 14},
        {"text": "Comandos confiáveis; telemetria periódica pode perder", "bullet": True, "size": 14},
        {"text": "Autenticação antes da concessão de controle", "bullet": True, "size": 14},
        {"text": "Relay primário + backup, replicação e failover", "bullet": True, "size": 14},
        {"text": "Heartbeat, detecção de falhas e preservação da posse", "bullet": True, "size": 14},
        {"text": "Execução local em processos ou threads, sempre por UDP real", "bullet": True, "size": 14},
    ])
    textbox(s, 6.85, 3.3, 5.7, .3, "FORA DO ESCOPO / LIMITAÇÕES", 14, True, RED)
    richtext(s, 6.83, 3.7, 5.8, 2.72, [
        {"text": "Relay geograficamente distribuído, consenso e eleição", "bullet": True, "size": 14},
        {"text": "QUIC padrão: TLS 1.3, congestion control, 0-RTT, migração", "bullet": True, "size": 14},
        {"text": "Streaming de vídeo real (apenas a necessidade é modelada)", "bullet": True, "size": 14},
        {"text": "Banco de identidade seguro; credenciais estão em memória", "bullet": True, "size": 14},
        {"text": "Criptografia/confidencialidade do tráfego de aplicação", "bullet": True, "size": 14},
        {"text": "Failback automático ao primário recuperado", "bullet": True, "size": 14},
    ])
    footer(s, "Código: relay_server.py, pilot_client.py, rov_simulator.py, run_demo.py e demo_dashboard.py.")

    # 3 — escolha de transporte
    s = add_slide(prs); title(s, 3, "Transporte: por que não ficar só com TCP ou UDP?",
                              "Delay é crítico para comando e vídeo; cada classe de mensagem pede uma garantia diferente.")
    headers = [("TCP", BLUE), ("UDP PURO", ORANGE), ("QUIC-LITE DO PROJETO", GREEN)]
    xs = [.65, 4.62, 8.59]
    bodies = [
        "✓ Entrega e ordem\n✓ Retransmissão\n✓ Conexão via 3-way handshake\n+ TLS exige handshake adicional\n\n✕ Head-of-line: uma perda bloqueia os bytes seguintes\n✕ Mais latência para dados perecíveis",
        "✓ Sem conexão\n✓ Cabeçalho pequeno\n✓ Baixa latência\n\n✕ Sem entrega garantida\n✕ Sem ordem\n✕ Sem autenticação/criptografia embutidas\n✕ Aplicação resolve o que precisar",
        "✓ UDP real\n✓ Canal confiável e ordenado\n✓ Canal “último valor vence”\n✓ Fluxos independentes\n\n✕ Não é QUIC padrão\n✕ Sem TLS/congestion control\n✕ Garantias didáticas, não industriais",
    ]
    for x, (h, c), b in zip(xs, headers, bodies):
        box(s, x, 1.48, 3.46, 4.58, h, b, LIGHT, c, c, BLACK)
    textbox(s, .8, 6.34, 11.7, .55,
            "DECISÃO  →  confiabilidade seletiva: proteger controle e autenticação sem fazer telemetria/vídeo antigo atrasar informação nova.",
            16, True, BLACK, PP_ALIGN.CENTER)
    footer(s, "Código: quiclite.py. Observação: QUIC real usa UDP, TLS 1.3 e múltiplos streams; esta implementação reproduz apenas parte das ideias.")

    # 4 — HOL
    s = add_slide(prs); title(s, 4, "Head-of-line blocking: a perda não pode parar tudo")
    textbox(s, .7, 1.4, 5.75, .32, "TCP · UM FLUXO ORDENADO", 14, True, BLUE)
    y = 2.05
    for i, (lab, col) in enumerate([("CMD 10", BLUE), ("TEL 11", GRAY), ("PERDIDO 12", RED), ("TEL 13", GRAY), ("CMD 14", BLUE)]):
        box(s, .7 + i*1.08, y, .94, .72, lab, "", WHITE, col, col)
    line(s, .85, 3.05, 5.75, 3.05, RED, 2)
    textbox(s, .8, 3.18, 5.2, .7, "Mesmo que 13 e 14 cheguem, a aplicação espera o pacote 12 ser retransmitido.", 15, False, RED)
    textbox(s, 6.85, 1.4, 5.75, .32, "QUIC-LITE · CANAIS INDEPENDENTES", 14, True, GREEN)
    pill(s, 6.9, 2.0, 1.45, "CONFIÁVEL", BLUE)
    box(s, 8.55, 1.9, 1.2, .72, "CMD 10", "", WHITE, BLUE, BLUE)
    box(s, 9.95, 1.9, 1.2, .72, "CMD 11", "", WHITE, BLUE, BLUE)
    line(s, 9.75, 2.26, 9.94, 2.26, BLUE, 2, True)
    pill(s, 6.9, 3.18, 1.45, "NÃO CONF.", GRAY)
    box(s, 8.55, 3.08, 1.2, .72, "TEL A", "", WHITE, GRAY, GRAY)
    box(s, 9.95, 3.08, 1.2, .72, "PERDEU", "", WHITE, RED, RED)
    box(s, 11.35, 3.08, 1.2, .72, "TEL C", "", WHITE, GRAY, GRAY)
    textbox(s, 6.95, 4.18, 5.45, .8,
            "A perda de telemetria não bloqueia comandos; o próximo estado substitui o anterior. No canal confiável, somente a sequência daquele peer aguarda.",
            15, False, GREEN)
    box(s, 1.25, 5.4, 10.8, .9, "IMPACTO OPERACIONAL",
        "Controle responsivo e streaming fluido importam mais que recuperar um frame ou uma leitura de sensor já obsoleta.",
        LIGHT, MID, BLACK, BLACK)
    footer(s, "Código: quiclite.Endpoint._recv_loop(), send_reliable(), send_unreliable() e _retransmit_loop().")

    # 5 — canais
    s = add_slide(prs); title(s, 5, "Dois canais sobre o mesmo socket UDP",
                              "A escolha acontece mensagem a mensagem.")
    box(s, .7, 1.42, 5.8, 4.82, "CANAL CONFIÁVEL + ORDENADO",
        "Mecanismo\nseq por peer · ACK · buffer fora de ordem · retransmissão após 0,4 s\n\nMensagens\nregister · registered · auth_challenge · auth_response · auth_ok/auth_fail · control_granted/denied · command · release_control · rov_offline · error · replicate\n\nSem limite de tentativas enquanto o peer existir.",
        rgb("EAF2FA"), BLUE, BLUE, BLACK)
    box(s, 6.83, 1.42, 5.8, 4.82, "CANAL NÃO CONFIÁVEL",
        "Mecanismo\nDatagrama único · seq = 0 · sem ACK · sem reenvio · entrega imediata\n\nMensagens\ntelemetry · heartbeat · relay_heartbeat · relay_ping\n\nPrincípio\nSão eventos periódicos: se um se perde, o próximo contém estado mais novo.",
        rgb("F3F4F5"), GRAY, GRAY, BLACK)
    textbox(s, 1.15, 6.5, 11.05, .42,
            "A confiabilidade é por endpoint remoto: o relay mantém estado separado para cada piloto, ROV e relay par.",
            15, True, BLACK, PP_ALIGN.CENTER)
    footer(s, "Código: protocol.py (classificação das mensagens) e quiclite.py (Peer/Endpoint).")

    # 6 — pacote
    s = add_slide(prs); title(s, 6, "Protocolo no fio: 5 bytes + JSON")
    textbox(s, .7, 1.45, 2.25, .3, "CABEÇALHO BINÁRIO", 14, True, BLUE)
    box(s, .7, 1.92, 2.0, 1.18, "BYTE 0", "tipo\n0 conf. · 1 não conf. · 2 ACK", rgb("EAF2FA"), BLUE, BLUE, BLACK)
    box(s, 2.72, 1.92, 3.25, 1.18, "BYTES 1…4", "sequência uint32\nbig-endian", rgb("EAF2FA"), BLUE, BLUE, BLACK)
    box(s, 5.99, 1.92, 6.63, 1.18, "BYTES 5…N", "payload JSON em UTF-8\nvazio para ACK", rgb("F3F4F5"), GRAY, GRAY, BLACK)
    textbox(s, .7, 3.45, 3.25, .3, "EXEMPLOS DE PAYLOAD", 14, True, GREEN)
    box(s, .7, 3.9, 3.75, 1.28, "COMANDO · CONFIÁVEL",
        '{"type":"command",\n "action":"thruster_frente","value":80}', WHITE, BLUE, BLUE, BLACK)
    box(s, 4.78, 3.9, 3.75, 1.28, "TELEMETRIA · NÃO CONF.",
        '{"type":"telemetry","battery":99.1,\n "depth":0.24,"temperature":18.2}', WHITE, GRAY, GRAY, BLACK)
    box(s, 8.86, 3.9, 3.75, 1.28, "ACK",
        "tipo = 2\nseq = pacote confirmado\npayload = vazio", WHITE, GREEN, GREEN, BLACK)
    richtext(s, .85, 5.58, 11.6, 1.0, [
        {"text": "Recepção confiável: confirma inclusive duplicatas; entrega somente a partir de recv_expected.", "bullet": True, "size": 14},
        {"text": "JSON inválido ou UTF-8 inválido é descartado; tamanho máximo recebido: 65.535 bytes.", "bullet": True, "size": 14},
    ])
    footer(s, "Código: quiclite.py — HEADER = struct.Struct('!BI'), _decode() e recvfrom(65535).")

    # 7 — auth
    s = add_slide(prs); title(s, 7, "Autenticação: desafio–resposta HMAC-SHA256",
                              "A senha não é enviada; o relay só concede controle após autenticação.")
    actors = [(1.05, "PILOTO", BLUE), (6.08, "RELAY", CYAN)]
    for x, a, c in actors:
        pill(s, x, 1.4, 2.1, a, c)
        line(s, x + 1.05, 1.78, x + 1.05, 6.3, c, 1.2, False, True)
    steps = [
        (2.05, 2.1, 6.0, "1  register(role=pilot, id, target)", BLUE, True),
        (2.48, 6.0, 3.15, "2  registered + auth_challenge(nonce 128 bits)", CYAN, False),
        (3.12, 2.1, 6.0, "3  auth_response = HMAC-SHA256(senha, nonce)", BLUE, True),
        (3.75, 6.0, 3.15, "4  compare_digest(esperado, recebido)", CYAN, False),
        (4.4, 6.0, 3.15, "5  auth_ok(token 128 bits) ou auth_fail", GREEN, False),
        (5.04, 6.0, 3.15, "6  control_granted / control_denied", GREEN, False),
    ]
    for y, x1, x2, txt, c, right in steps:
        line(s, x1, y, x2, y, c, 1.8, True)
        textbox(s, min(x1, x2)+.2, y-.28, abs(x2-x1)-.4, .26, txt, 11.5, True, c,
                PP_ALIGN.CENTER)
    box(s, 9.45, 1.42, 3.1, 4.85, "PROPRIEDADES",
        "• Nonce aleatório por registro\n\n• Replay de resposta antiga falha\n\n• Comparação em tempo constante\n\n• Piloto desconhecido é recusado\n\n• Token opaco identifica a sessão\n\nAtenção: o token é emitido e armazenado, mas não é exigido nos comandos.",
        LIGHT, MID, ORANGE, BLACK)
    footer(s, "Código: protocol.py e relay_server.py::_handle_register/_handle_auth; pilot_client.py::_on_message.")

    # 8 — credenciais
    s = add_slide(prs); title(s, 8, "“Banco” de credenciais: adequado apenas à demonstração")
    box(s, .8, 1.45, 5.25, 3.65, "IMPLEMENTAÇÃO ATUAL",
        'PILOT_CREDENTIALS = {\n  "pilotoA": "mergulho2026",\n  "pilotoB": "trocaraki"\n}\n\nOs dois relays carregam a mesma tabela para poder autenticar após o failover.',
        rgb("FFF4E8"), ORANGE, ORANGE, BLACK)
    box(s, 6.55, 1.45, 5.95, 3.65, "SISTEMA REAL",
        "• Serviço de identidade ou banco protegido\n• Senhas com hash lento + salt (Argon2id/bcrypt/scrypt)\n• Segredos fora do código e com rotação\n• Controle de acesso, auditoria e rate limiting\n• Replicação segura/consistente das identidades\n• TLS/mTLS para confidencialidade e identidade do relay",
        rgb("EEF7F1"), GREEN, GREEN, BLACK)
    textbox(s, 1.05, 5.55, 11.3, .85,
            "O HMAC evita transmitir a senha, mas não corrige senha em texto puro no código, senha fraca, ataque offline após captura do desafio/resposta ou ausência de canal autenticado.",
            17, True, RED, PP_ALIGN.CENTER)
    footer(s, "Código: protocol.py::PILOT_CREDENTIALS. Credenciais incluídas aqui apenas para documentar fielmente a demo.")

    # 9 — failover
    s = add_slide(prs); title(s, 9, "Detecção de falhas e failover do relay",
                              "Os timeouts foram ordenados para o backup assumir antes da migração dos clientes.")
    labels = [
        (1.0, "0 s", "relay_ping / relay_heartbeat\ncada 1 s", GREEN),
        (3.25, "2,5 s", "backup detecta silêncio\ne torna-se ATIVO", ORANGE),
        (6.1, "6 s", "piloto/ROV detectam silêncio\ne migram", BLUE),
        (9.05, "≤ 12 s", "reserva preserva o dono\nanterior do ROV", CYAN),
    ]
    line(s, 1.35, 2.5, 11.75, 2.5, BLACK, 2.2, True)
    for x, t, b, c in labels:
        sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.21), Inches(.58), Inches(.58))
        sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
        textbox(s, x-.12, 1.43, .82, .3, t, 14, True, c, PP_ALIGN.CENTER)
        textbox(s, x-.4, 2.95, 1.4, .78, b, 12.5, False, BLACK, PP_ALIGN.CENTER)
    box(s, .75, 4.18, 3.7, 1.55, "CLIENTES",
        "heartbeat a cada 1 s; relay remove cliente após 6 s sem mensagem. Qualquer mensagem atualiza last_seen.",
        LIGHT, MID, BLUE, BLACK)
    box(s, 4.82, 4.18, 3.7, 1.55, "REPLICAÇÃO",
        "Primário envia rov_up/down, pilot_up/down e control pelo canal confiável. Backup mantém estado-espelho.",
        LIGHT, MID, CYAN, BLACK)
    box(s, 8.89, 4.18, 3.7, 1.55, "RECUPERAÇÃO",
        "Clientes removem o peer morto, alternam relay e se registram de novo; o piloto refaz a autenticação.",
        LIGHT, MID, GREEN, BLACK)
    textbox(s, .8, 6.35, 11.75, .42,
            "Sem consenso/split-brain protection: é um esquema primário–backup didático, não uma implantação distribuída completa.",
            15, True, RED, PP_ALIGN.CENTER)
    footer(s, "Código: relay_server.py (timeouts, replicação, reserva); pilot_client.py e rov_simulator.py (failover).")

    # 10 — segurança
    s = add_slide(prs); title(s, 10, "Segurança e robustez: o que existe — e o que falta")
    box(s, .72, 1.4, 3.75, 4.95, "CONTROLES IMPLEMENTADOS",
        "✓ HMAC-SHA256 desafio–resposta\n✓ nonce aleatório de 128 bits\n✓ token de sessão aleatório\n✓ compare_digest contra timing\n✓ autenticação antes do comando\n✓ exclusão mútua por ROV\n✓ duplicatas confiáveis não são reaplicadas\n✓ JSON corrompido é ignorado\n✓ liveness + failover + reserva de posse",
        rgb("EEF7F1"), GREEN, GREEN, BLACK)
    box(s, 4.8, 1.4, 3.75, 4.95, "ATAQUES MITIGADOS",
        "• Replay simples da autenticação\n• Senha em trânsito\n• Piloto não cadastrado\n• Corrida normal entre pilotos\n• Perda, duplicação e reordenação de datagramas confiáveis\n• Falha silenciosa de clientes\n• Queda do relay primário",
        rgb("EAF2FA"), BLUE, BLUE, BLACK)
    box(s, 8.88, 1.4, 3.75, 4.95, "LACUNAS PARA PRODUÇÃO",
        "✕ tráfego sem criptografia/TLS\n✕ relay não autenticado\n✕ comandos não carregam/validam token\n✕ sem integridade HMAC por pacote\n✕ sem expiração/revogação de token\n✕ sem rate limiting / anti-DoS\n✕ sem limite de buffers e retransmissões\n✕ credenciais em texto puro\n✕ sem consenso contra split-brain",
        rgb("FCEEEE"), RED, RED, BLACK)
    footer(s, "Análise baseada em protocol.py, quiclite.py e relay_server.py. As lacunas são inferências diretas do código.")

    # 11 — Tkinter
    s = add_slide(prs); title(s, 11, "Tkinter: visualização sem alterar a rede",
                              "A lógica é separada da interface; os testes rodam headless.")
    box(s, .75, 1.42, 3.55, 4.78, "MULTI-JANELA",
        "run_demo.py\n\nCada nó em um processo e uma janela nativa. Reforça a separação entre hosts e mantém comunicação por UDP.",
        LIGHT, MID, BLUE, BLACK)
    box(s, 4.88, 1.42, 3.55, 4.78, "DASHBOARD",
        "demo_dashboard.py\n\nMesmos nós como threads no mesmo processo, ainda via sockets UDP no loopback. Mostra topologia, pacotes, perda e cena submarina.",
        LIGHT, MID, CYAN, BLACK)
    box(s, 9.0, 1.42, 3.55, 4.78, "THREAD-SAFETY",
        "gui_common.py\n\nThreads de rede colocam eventos em queue.Queue; a thread principal drena a fila com .after(80 ms). Widgets nunca são atualizados diretamente pela rede.",
        LIGHT, MID, GREEN, BLACK)
    pill(s, 1.25, 5.65, 2.55, "PROCESSOS", BLUE)
    line(s, 3.8, 5.82, 5.22, 5.82, BLACK, 1.8, True)
    pill(s, 5.25, 5.65, 2.55, "UDP / LOOPBACK", CYAN)
    line(s, 7.8, 5.82, 9.22, 5.82, BLACK, 1.8, True)
    pill(s, 9.25, 5.65, 2.55, "EVENTOS → UI", GREEN)
    textbox(s, 1.0, 6.52, 11.35, .35,
            "Demonstração: conectar piloto · comandar · adicionar concorrente · simular perda · derrubar primário.",
            15, True, BLACK, PP_ALIGN.CENTER)
    footer(s, "Código: gui_common.py, run_demo.py, demo_dashboard.py e classes *Node sem dependência de UI.")

    # 12 — diagrama final
    s = add_slide(prs); title(s, 12, "Visão completa do sistema",
                              "Fluxo normal, canais, autenticação, replicação e failover em uma única vista.")
    # nodes
    box(s, .55, 2.02, 2.3, 1.35, "PILOTO A",
        "auth ✓\ncontrola rov1", rgb("EAF2FA"), BLUE, BLUE, BLACK)
    box(s, .55, 4.48, 2.3, 1.35, "PILOTO B",
        "auth ✓\ncontrole negado", rgb("EAF2FA"), BLUE, BLUE, BLACK)
    box(s, 5.0, 1.42, 3.0, 1.52, "RELAY PRIMÁRIO",
        ":5000 · ativo\nsessões + arbitragem", rgb("EEF7F1"), GREEN, GREEN, BLACK)
    box(s, 5.0, 4.92, 3.0, 1.52, "RELAY BACKUP",
        ":5001 · espelho\nassume após 2,5 s", rgb("FFF4E8"), ORANGE, ORANGE, BLACK)
    box(s, 10.05, 2.84, 2.72, 1.75, "ROV 1",
        "sensores + thrusters\ntelemetria a cada 1,5 s", rgb("EEF7F1"), CYAN, CYAN, BLACK)
    # links
    line(s, 2.85, 2.48, 5.0, 2.05, BLUE, 2.2, True)
    textbox(s, 2.9, 1.72, 1.9, .48, "registro · HMAC\ncomandos [CONF.]", 10.5, True, BLUE, PP_ALIGN.CENTER)
    line(s, 2.85, 5.08, 5.0, 2.58, BLUE, 1.5, True, True)
    textbox(s, 3.25, 4.05, 1.55, .5, "autentica\nmas não obtém posse", 10, False, BLUE, PP_ALIGN.CENTER)
    line(s, 8.0, 2.15, 10.05, 3.33, BLUE, 2.2, True)
    textbox(s, 8.1, 2.25, 1.72, .45, "comando\n[CONF.]", 10.5, True, BLUE, PP_ALIGN.CENTER)
    line(s, 10.05, 4.05, 8.0, 2.58, GRAY, 2.2, True)
    textbox(s, 8.12, 3.68, 1.7, .45, "telemetria\n[NÃO CONF.]", 10.5, True, GRAY, PP_ALIGN.CENTER)
    line(s, 6.5, 2.94, 6.5, 4.92, CYAN, 2.2, True)
    textbox(s, 6.73, 3.48, 1.85, .48, "replicate [CONF.]\nrelay_ping [NÃO CONF.]", 10.5, True, CYAN)
    line(s, 8.0, 5.68, 10.05, 4.25, ORANGE, 1.6, True, True)
    line(s, 5.0, 5.68, 2.85, 5.5, ORANGE, 1.6, True, True)
    textbox(s, 8.35, 5.05, 1.35, .52, "failover\n+ re-registro", 10.5, True, ORANGE, PP_ALIGN.CENTER)
    # legend
    pill(s, .7, 6.58, 2.05, "CONF. · ACK + RETX", BLUE)
    pill(s, 3.0, 6.58, 2.35, "NÃO CONF. · ÚLTIMO VENCE", GRAY)
    pill(s, 5.6, 6.58, 2.08, "HMAC + NONCE", GREEN)
    pill(s, 7.93, 6.58, 2.1, "FAILOVER · 6 s", ORANGE)
    pill(s, 10.28, 6.58, 2.18, "RESERVA · 12 s", CYAN)
    footer(s, "Síntese do código: quiclite.py, protocol.py, relay_server.py, pilot_client.py e rov_simulator.py.")

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
